"""Generate the submission slides.pptx (python-pptx) - ~12 concise slides for the imgtrans
document-image-translation system. Degrades to a Markdown outline if python-pptx is unavailable.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from ..config import AppConfig, run_dir
from ..logging_utils import get_logger
from . import charts as charts_mod
from .artifact_loader import e2e_metric, headline, load_artifacts, model_version, mt_metric, ocr_metric

logger = get_logger(__name__)


def _slides(cfg: AppConfig, arts: Dict[str, Any]) -> List[Tuple[str, List[str]]]:
    e2e = e2e_metric(arts, "chrf")
    cer = ocr_metric(arts, "cer")
    fit = headline(arts, "mean_fit_rate")
    res = (f"end-to-end image-translation chrF {e2e:.1f}" if e2e is not None
           else "train + evaluate to populate results")
    ocrline = (f"OCR CER {cer:.3f}; overlay fit-rate {fit:.2f}"
               if (cer is not None and fit is not None) else "OCR CER/WER + overlay fit-rate")
    d = f"{cfg.mt.src_lang} -> {cfg.mt.tgt_lang}"
    return [
        ("Document-Image Machine Translation",
         [f"{cfg.author} - Student {cfg.student_id}", "NLP in Industry - Final Assignment",
          f"Translate the text INSIDE an image / scan ({d}) and render it back",
          "Cascade: OCR -> MT -> layout-preserving overlay",
          "Low-confidence / poor-fit output is flagged for human review"]),
        ("Business Problem & Motivation",
         ["Travelers, back-office, localization need to translate text in photos/scans",
          "More than 'OCR then translate': keep the layout (overlay) + gate quality",
          "The fix = an OCR front-end + a trainable MT core + an overlay renderer + a QA agent",
          "Only the MT model is trained (the NLP heart); OCR + render are algorithmic"]),
        ("Proposed Solution",
         ["OCR the page into text blocks with boxes + confidence (Tesseract)",
          "Translate each block with a fine-tuned m2m100 (the trainable core)",
          "Render the translation back into the original box (auto font-fit + wrap)",
          "A 5-decision agent gates OCR confidence, translation quality, overlay fit"]),
        ("System Architecture",
         ["ingest (D1 route + quality) -> ocr (D2 born-digital vs scanned)",
          "-> translate (D3 OCR-confidence gate) -> verify (D4 round-trip)",
          "-> render (D5 fit gate: overlay | side-by-side | needs_review)",
          "Runs fully offline (SeedEngine OCR + dictionary MT) for tests/CI"]),
        ("Data & The Synthetic Generator",
         ["NO real in-image translation benchmark exists -> synthetic generator",
          "Render source sentences onto pages + embed gold (source, translation, boxes)",
          "MT fine-tune: Helsinki-NLP/opus-100 en-fr (license unknown flag)",
          f"Default {d}; offline seed = synthetic pages + en->fr dictionary"]),
        ("The Trainable MT Core",
         ["Base: facebook/m2m100_418M (MIT, many-to-many ~100 langs)",
          "HF Seq2SeqTrainer, chrF selection, bf16/tf32 on H100, early stopping",
          "Baselines: zero-shot MT + dictionary + identity floor",
          "Alts: opus-mt-en-fr (Apache, T4), mbart-50 (MIT); NLLB (CC-BY-NC, flagged)"]),
        ("OCR Front-End & Overlay Renderer",
         ["Tesseract (Apache) gives word boxes + confidence; PyMuPDF born-digital router",
          "PIL-only fit-to-box: binary-search font size + greedy word wrap",
          "Value-add: re-fit the longer translation into the source box",
          "Poor fit -> side-by-side fallback; low OCR conf -> skip the block"]),
        ("Agentic AI Component (D1-D5)",
         ["Deterministic FSM + optional LLM brain (OFF by default)",
          "D1 input router+quality - D2 born-digital vs scanned",
          "D3 per-block OCR-confidence gate - D4 round-trip back-translation verify",
          "D5 render-fit feasibility gate (overlay | side-by-side | needs_review)"]),
        ("Evaluation Results",
         [res, ocrline,
          "MT chrF/BLEU + OCR CER/WER + end-to-end chrF + layout fit-rate reported jointly",
          "Offline floor saturates; the fine-tuned m2m100 dominates on real opus-100 pairs"]),
        ("Deployment Overview",
         ["FastAPI /translate-text + /translate-image (multipart, base64 overlay) + /healthz",
          "Gradio demo (upload image -> translated image + text)",
          "Docker (tesseract-ocr + libGL) + HF Space; lazy deps + offline fallback",
          "Metadata-only job logging"]),
        ("Continual Learning, Monitoring & Ethics",
         ["Collect corrected pages -> re-fine-tune MT -> promote if chrF non-regressing",
          "monitor-log: needs-review rate + OCR-conf/fit-rate drift + latency",
          "Privacy: document images = sensitive PII -> local processing, no raw retention",
          "Mistranslation harm -> verify gate + abstention; the tool ASSISTS, flags for review"]),
        ("Key Takeaways & Future Work",
         ["A layout-preserving, confidence-gated, debuggable image-translation pipeline",
          "The overlay renderer + the OCR-confidence/verify gates are the value-add",
          "Future: neural OCR (TrOCR/GOT-OCR2), inpainting, multi-page overlay, more languages",
          "Future: end-to-end image-to-image MT + terminology memory"]),
    ]


def generate_slides(cfg: AppConfig, title: Optional[str] = None, author: Optional[str] = None,
                    out_path: Optional[str] = None) -> str:
    arts = load_artifacts(cfg)
    out_path = Path(out_path) if out_path else run_dir() / "report" / "slides.pptx"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    slides = _slides(cfg, arts)
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except Exception as exc:
        logger.warning("python-pptx unavailable (%s); writing markdown outline", exc)
        md = "\n\n".join(f"## {t}\n" + "\n".join(f"- {b}" for b in bs) for t, bs in slides)
        alt = out_path.with_suffix(".md")
        alt.write_text(md, encoding="utf-8")
        return str(alt)

    try:
        chart = charts_mod.e2e_chart(arts, run_dir() / "report" / "slide_e2e.png")
    except Exception:
        chart = None
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    accent = RGBColor(0x2B, 0x6C, 0xB0)
    for i, (t, bullets) in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
        tf = bar.text_frame; tf.text = t
        tf.paragraphs[0].font.size = Pt(28); tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.5),
                                        Inches(8.3 if (i == 8 and chart) else 12), Inches(5.4))
        bt = body.text_frame; bt.word_wrap = True
        for j, bp in enumerate(bullets):
            p = bt.paragraphs[0] if j == 0 else bt.add_paragraph()
            p.text = "-  " + bp; p.font.size = Pt(20); p.space_after = Pt(10)
        if i == 8 and chart:
            slide.shapes.add_picture(str(chart), Inches(8.9), Inches(1.7), width=Inches(4.0))
        foot = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4))
        foot.text_frame.text = f"{title or cfg.project_title} - {author or cfg.author} ({cfg.student_id})"
        foot.text_frame.paragraphs[0].font.size = Pt(9)
    prs.save(str(out_path))
    logger.info("Slides -> %s", out_path)
    return str(out_path)


__all__ = ["generate_slides"]
